"""Extract original experimental image blobs from the Figure 3 source slide.

The PowerPoint slide is used only as an experimental-image container. No
bioinformatic plot is cropped from the slide. Output names are stable so the
assembly script can link them reproducibly.
"""

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SOURCE = Path(
    r"F:\20250325cell背靠背拒稿\20251214上皮化\20260313重新拼图"
    r"\20260621-上皮化文章拼图.pptx"
)
OUTPUT = Path(__file__).resolve().parents[1] / "assets" / "linked_rasters"
SLIDE_NUMBER = 5


def picture_at_path(shapes, path: str):
    current = shapes
    shape = None
    for token in path.split("."):
        shape = current[int(token)]
        current = shape.shapes if shape.shape_type == MSO_SHAPE_TYPE.GROUP else None
    if shape is None or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise ValueError(f"Shape path {path} is not a picture")
    return shape


def main():
    OUTPUT.mkdir(parents=True, exist_ok=True)
    slide = Presentation(str(SOURCE)).slides[SLIDE_NUMBER - 1]
    selections = {
        # VIM/KRT5 wound-tissue image and its H&E inset.
        "Figure3_H_VIM_KRT5.png": "22.0",
        "Figure3_H_HE_inset.png": "22.1",
        # Pdgfra lineage-tracing composite.
        "Figure3_I_Pdgfra_lineage.png": "14",
        # Sorted-cell IF composite including the source quantification panel.
        "Figure3_J_sorted_cell_IF.png": "17",
    }
    for filename, shape_path in selections.items():
        shape = picture_at_path(slide.shapes, shape_path)
        target = OUTPUT / filename
        target.write_bytes(shape.image.blob)
        print(f"{filename}\t{shape.image.size[0]}x{shape.image.size[1]}\t{target}")


if __name__ == "__main__":
    main()
